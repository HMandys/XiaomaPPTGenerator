#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
试卷分析PPT生成器
功能：导入成绩单、上传试卷、拉框切割、生成PPT
author: HMandy
大部分由Claude 生成，仅供参考
"""

import sys
import json
import os
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout,
                             QHBoxLayout, QPushButton, QLabel, QFileDialog,
                             QMessageBox, QScrollArea, QLineEdit, QSpinBox,
                             QComboBox, QGroupBox)
from PyQt5.QtGui import QPixmap, QPainter, QPen, QImage
from PyQt5.QtCore import Qt, QRect, QPoint
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import io

try:
    import fitz  # PyMuPDF

    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("警告: PyMuPDF未安装，PDF功能不可用")


class ImageCropWidget(QWidget):
    """图片切割控件，支持拉框选择题目区域"""

    def __init__(self, image_path, is_pdf=False, pdf_page=0):
        super().__init__()
        self.image_path = image_path
        self.is_pdf = is_pdf
        self.pdf_page = pdf_page
        self.pdf_doc = None
        self.original_image = None

        # 加载图片
        if is_pdf and PDF_SUPPORT:
            self.pdf_doc = fitz.open(image_path)
            page = self.pdf_doc[pdf_page]

            # 将PDF页面转换为图片 (提高分辨率)
            mat = fitz.Matrix(2.0, 2.0)  # 2倍缩放以提高清晰度
            pix = page.get_pixmap(matrix=mat)

            # 转换为QPixmap
            img_data = pix.tobytes("png")
            qimage = QImage.fromData(img_data)
            self.original_pixmap = QPixmap.fromImage(qimage)

            # 保存原始PIL Image用于后续切割
            self.original_image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        else:
            self.original_pixmap = QPixmap(image_path)
            self.original_image = Image.open(image_path)

        self.display_pixmap = self.original_pixmap.copy()
        self.scale_factor = 1.0

        # 计算缩放比例以适应屏幕
        screen_size = QApplication.desktop().screenGeometry()
        max_width = int(screen_size.width() * 0.8)
        max_height = int(screen_size.height() * 0.7)

        if self.original_pixmap.width() > max_width or self.original_pixmap.height() > max_height:
            self.display_pixmap = self.original_pixmap.scaled(
                max_width, max_height, Qt.KeepAspectRatio, Qt.SmoothTransformation
            )
            self.scale_factor = self.display_pixmap.width() / self.original_pixmap.width()

        self.setFixedSize(self.display_pixmap.size())

        self.start_point = None
        self.end_point = None
        self.rectangles = []  # 存储所有矩形区域 [(rect, question_number), ...]
        self.current_question_number = 1
        self.is_drawing = False

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.drawPixmap(0, 0, self.display_pixmap)

        # 绘制已保存的矩形
        pen = QPen(Qt.green, 2, Qt.SolidLine)
        painter.setPen(pen)
        for rect, qnum in self.rectangles:
            painter.drawRect(rect)
            # 在矩形上方显示题号
            painter.drawText(rect.topLeft() + QPoint(5, -5), f"题{qnum}")

        # 绘制当前正在画的矩形
        if self.start_point and self.end_point:
            pen = QPen(Qt.red, 2, Qt.DashLine)
            painter.setPen(pen)
            rect = QRect(self.start_point, self.end_point).normalized()
            painter.drawRect(rect)

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.start_point = event.pos()
            self.is_drawing = True

    def mouseMoveEvent(self, event):
        if self.is_drawing:
            self.end_point = event.pos()
            self.update()

    def mouseReleaseEvent(self, event):
        if event.button() == Qt.LeftButton and self.is_drawing:
            self.end_point = event.pos()
            self.is_drawing = False

            # 保存矩形区域
            rect = QRect(self.start_point, self.end_point).normalized()
            if rect.width() > 10 and rect.height() > 10:  # 忽略太小的矩形
                self.rectangles.append((rect, self.current_question_number))
                self.current_question_number += 1
                self.update()
                # 通知父窗口检查是否可以生成PPT
                if self.parent() and hasattr(self.parent().parent(), 'check_ready'):
                    self.parent().parent().check_ready()

            self.start_point = None
            self.end_point = None

    def get_cropped_regions(self):
        """获取所有切割区域的图片"""
        regions = []

        # 使用已加载的原始图片
        original_image = self.original_image

        # 按题号排序
        sorted_rects = sorted(self.rectangles, key=lambda x: x[1])

        for rect, qnum in sorted_rects:
            # 将显示坐标转换为原始图片坐标
            x1 = int(rect.x() / self.scale_factor)
            y1 = int(rect.y() / self.scale_factor)
            x2 = int(rect.right() / self.scale_factor)
            y2 = int(rect.bottom() / self.scale_factor)

            # 切割图片
            cropped = original_image.crop((x1, y1, x2, y2))
            regions.append((qnum, cropped))

        return regions

    def __del__(self):
        """清理资源"""
        if self.pdf_doc:
            self.pdf_doc.close()

    def clear_last(self):
        """删除最后一个矩形"""
        if self.rectangles:
            self.rectangles.pop()
            self.current_question_number -= 1
            self.update()

    def clear_all(self):
        """清除所有矩形"""
        self.rectangles = []
        self.current_question_number = 1
        self.update()

    def save_regions_config(self, filename):
        """保存切割配置到文件"""
        config = []
        for rect, qnum in self.rectangles:
            config.append({
                'question_number': qnum,
                'x': rect.x(),
                'y': rect.y(),
                'width': rect.width(),
                'height': rect.height()
            })

        with open(filename, 'w', encoding='utf-8') as f:
            json.dump(config, f, ensure_ascii=False, indent=2)

    def load_regions_config(self, filename):
        """从文件加载切割配置"""
        try:
            with open(filename, 'r', encoding='utf-8') as f:
                config = json.load(f)

            self.rectangles = []
            max_qnum = 0
            for item in config:
                rect = QRect(
                    int(item['x']), int(item['y']),
                    int(item['width']), int(item['height'])
                )
                qnum = item['question_number']
                self.rectangles.append((rect, qnum))
                max_qnum = max(max_qnum, qnum)

            self.current_question_number = max_qnum + 1
            self.update()
            return True
        except Exception as e:
            print(f"加载配置失败: {e}")
            return False


class MainWindow(QMainWindow):
    """主窗口"""

    def __init__(self):
        super().__init__()
        self.setWindowTitle("试卷分析PPT生成器")
        self.setGeometry(100, 100, 1200, 800)

        self.excel_file = None
        self.exam_image = None
        self.crop_widget = None
        self.df = None
        self.is_pdf = False
        self.pdf_total_pages = 0
        self.current_pdf_page = 0

        self.init_ui()

    def init_ui(self):
        """初始化UI"""
        central_widget = QWidget()
        self.setCentralWidget(central_widget)

        main_layout = QVBoxLayout()
        central_widget.setLayout(main_layout)

        # 文件选择区域
        file_layout = QHBoxLayout()

        # Excel文件选择
        self.excel_label = QLabel("成绩单：未选择")
        self.excel_btn = QPushButton("选择成绩单Excel")
        self.excel_btn.clicked.connect(self.load_excel)
        file_layout.addWidget(self.excel_label)
        file_layout.addWidget(self.excel_btn)

        # 试卷图片/PDF选择
        self.image_label = QLabel("试卷：未选择")
        self.image_btn = QPushButton("选择试卷图片/PDF")
        self.image_btn.clicked.connect(self.load_image)
        file_layout.addWidget(self.image_label)
        file_layout.addWidget(self.image_btn)

        main_layout.addLayout(file_layout)

        # PDF页面选择区域
        self.pdf_page_widget = QWidget()
        pdf_page_layout = QHBoxLayout()
        pdf_page_layout.addWidget(QLabel("PDF页面："))
        self.pdf_page_combo = QComboBox()
        self.pdf_page_combo.currentIndexChanged.connect(self.change_pdf_page)
        pdf_page_layout.addWidget(self.pdf_page_combo)
        pdf_page_layout.addStretch()
        self.pdf_page_widget.setLayout(pdf_page_layout)
        self.pdf_page_widget.setVisible(False)
        main_layout.addWidget(self.pdf_page_widget)

        # 图片显示和切割区域
        self.scroll_area = QScrollArea()
        self.scroll_area.setWidgetResizable(True)
        main_layout.addWidget(self.scroll_area)

        # 控制按钮区域
        control_layout = QHBoxLayout()

        self.clear_last_btn = QPushButton("撤销上一个")
        self.clear_last_btn.clicked.connect(self.clear_last_rect)
        self.clear_last_btn.setEnabled(False)
        control_layout.addWidget(self.clear_last_btn)

        self.clear_all_btn = QPushButton("清除所有")
        self.clear_all_btn.clicked.connect(self.clear_all_rects)
        self.clear_all_btn.setEnabled(False)
        control_layout.addWidget(self.clear_all_btn)

        self.save_config_btn = QPushButton("保存切割配置")
        self.save_config_btn.clicked.connect(self.save_config)
        self.save_config_btn.setEnabled(False)
        control_layout.addWidget(self.save_config_btn)

        self.load_config_btn = QPushButton("加载切割配置")
        self.load_config_btn.clicked.connect(self.load_config)
        self.load_config_btn.setEnabled(False)
        control_layout.addWidget(self.load_config_btn)

        control_layout.addStretch()

        self.generate_btn = QPushButton("生成PPT")
        self.generate_btn.clicked.connect(self.generate_ppt)
        # 始终启用按钮，让用户可以点击并获得反馈
        self.generate_btn.setEnabled(True)
        self.generate_btn.setStyleSheet("background-color: #4CAF50; color: white; font-size: 14px; padding: 10px;")
        control_layout.addWidget(self.generate_btn)

        # 添加调试按钮
        self.debug_btn = QPushButton("【调试】强制生成")
        self.debug_btn.clicked.connect(self.debug_generate)
        self.debug_btn.setStyleSheet("background-color: #FF9800; color: white; font-size: 12px; padding: 5px;")
        control_layout.addWidget(self.debug_btn)

        main_layout.addLayout(control_layout)

        # 说明文字
        help_text = QLabel("使用说明：\n"
                           "1. 选择成绩单Excel文件（需包含学生姓名和各题得分列）\n"
                           "2. 选择试卷图片文件或PDF文件\n"
                           "3. 如果是PDF，可以选择要分析的页面\n"
                           "4. 在试卷图片上用鼠标拉框选择每道题的区域（按顺序）\n"
                           "5. 点击'生成PPT'创建分析报告")
        help_text.setStyleSheet("color: #666; padding: 10px;")
        main_layout.addWidget(help_text)

    def load_excel(self):
        """加载Excel文件"""
        filename, _ = QFileDialog.getOpenFileName(
            self, "选择成绩单Excel文件", "", "Excel Files (*.xlsx *.xls)"
        )

        if filename:
            try:
                print(f"正在加载Excel文件: {filename}")

                # 先读取原始数据查看结构
                df_raw = pd.read_excel(filename, header=None)
                print(f"原始数据形状: {df_raw.shape}")
                print(f"前3行数据:")
                print(df_raw.head(3))

                # 查找真正的表头行
                header_row = None
                for i in range(min(5, len(df_raw))):  # 检查前5行
                    row_values = df_raw.iloc[i].astype(str).tolist()
                    print(f"第{i}行: {row_values[:5]}...")  # 显示前5个值
                    if any('姓名' in str(val) for val in row_values):
                        header_row = i
                        print(f"找到表头在第{i}行")
                        break

                if header_row is not None:
                    # 使用找到的行作为表头
                    self.df = pd.read_excel(filename, header=header_row)
                    print(f"使用第{header_row}行作为表头")
                    print(f"列名: {self.df.columns.tolist()[:10]}...")

                    # 删除"得分"标签行（如果存在）
                    if len(self.df) > 0:
                        first_data_row = self.df.iloc[0].astype(str).tolist()
                        if '得分' in ' '.join(first_data_row):
                            print("删除'得分'标签行")
                            self.df = self.df.iloc[1:].reset_index(drop=True)
                else:
                    # 如果没找到表头，使用默认方式
                    self.df = pd.read_excel(filename)
                    print("使用默认方式加载")

                # 清理空行
                self.df = self.df.dropna(how='all').reset_index(drop=True)

                self.excel_file = filename
                self.excel_label.setText(f"成绩单：{os.path.basename(filename)}")

                # 显示数据预览
                student_count = len(self.df)
                question_cols = [col for col in self.df.columns if '满分' in str(col) or '题' in str(col)]
                question_count = len(question_cols)

                print(f"学生数: {student_count}")
                print(f"题目列: {question_cols[:5]}...")  # 显示前5个题目列
                print(f"题目数: {question_count}")

                # 数据校验
                validation_result = self.validate_data()

                msg = f"已加载成绩单\n"
                msg += f"学生数: {student_count}\n"
                msg += f"检测到题目数: {question_count}\n\n"
                msg += "=== 数据校验 ===\n"
                msg += validation_result

                QMessageBox.information(self, "成功", msg)
                self.check_ready()
            except Exception as e:
                print(f"加载Excel失败: {e}")
                QMessageBox.critical(self, "错误", f"加载Excel文件失败：\n{str(e)}")
                import traceback
                traceback.print_exc()

    def load_image(self):
        """加载试卷图片或PDF"""
        # 根据是否支持PDF调整文件过滤器
        if PDF_SUPPORT:
            file_filter = "图片和PDF文件 (*.png *.jpg *.jpeg *.bmp *.pdf);;图片文件 (*.png *.jpg *.jpeg *.bmp);;PDF文件 (*.pdf);;所有文件 (*.*)"
        else:
            file_filter = "图片文件 (*.png *.jpg *.jpeg *.bmp);;所有文件 (*.*)"

        filename, _ = QFileDialog.getOpenFileName(
            self, "选择试卷图片或PDF", "", file_filter
        )

        if filename:
            try:
                self.exam_image = filename
                file_ext = os.path.splitext(filename)[1].lower()

                # 检查是否为PDF
                if file_ext == '.pdf':
                    if not PDF_SUPPORT:
                        QMessageBox.critical(self, "错误",
                                             "PDF支持未安装！\n"
                                             "请运行: pip install PyMuPDF")
                        return

                    self.is_pdf = True
                    # 获取PDF页数
                    pdf_doc = fitz.open(filename)
                    self.pdf_total_pages = len(pdf_doc)
                    pdf_doc.close()

                    # 显示PDF页面选择
                    self.pdf_page_combo.clear()
                    for i in range(self.pdf_total_pages):
                        self.pdf_page_combo.addItem(f"第 {i + 1} 页")
                    self.pdf_page_widget.setVisible(True)
                    self.current_pdf_page = 0

                    self.image_label.setText(f"试卷：{os.path.basename(filename)} (PDF, {self.pdf_total_pages}页)")
                else:
                    self.is_pdf = False
                    self.pdf_page_widget.setVisible(False)
                    self.image_label.setText(f"试卷：{os.path.basename(filename)}")

                # 创建切割控件
                self.crop_widget = ImageCropWidget(filename, self.is_pdf, self.current_pdf_page)
                self.scroll_area.setWidget(self.crop_widget)

                self.clear_last_btn.setEnabled(True)
                self.clear_all_btn.setEnabled(True)
                self.save_config_btn.setEnabled(True)
                self.load_config_btn.setEnabled(True)

                msg = "请用鼠标在试卷上拉框选择每道题的区域\n按照题号顺序依次框选"
                if self.is_pdf:
                    msg += f"\n\n当前显示：第{self.current_pdf_page + 1}页（共{self.pdf_total_pages}页）"

                QMessageBox.information(self, "提示", msg)
                self.check_ready()
            except Exception as e:
                QMessageBox.critical(self, "错误", f"加载文件失败：\n{str(e)}")
                import traceback
                traceback.print_exc()

    def change_pdf_page(self, index):
        """切换PDF页面"""
        if self.is_pdf and self.exam_image:
            try:
                self.current_pdf_page = index

                # 重新创建切割控件
                self.crop_widget = ImageCropWidget(self.exam_image, self.is_pdf, self.current_pdf_page)
                self.scroll_area.setWidget(self.crop_widget)

                QMessageBox.information(self, "提示",
                                        f"已切换到第 {index + 1} 页\n"
                                        "之前的框选已清除，请重新框选")
            except Exception as e:
                QMessageBox.critical(self, "错误", f"切换页面失败：\n{str(e)}")

    def clear_last_rect(self):
        """清除最后一个矩形"""
        if self.crop_widget:
            self.crop_widget.clear_last()

    def clear_all_rects(self):
        """清除所有矩形"""
        if self.crop_widget:
            reply = QMessageBox.question(self, "确认", "确定要清除所有选区吗？",
                                         QMessageBox.Yes | QMessageBox.No)
            if reply == QMessageBox.Yes:
                self.crop_widget.clear_all()

    def save_config(self):
        """保存切割配置"""
        if self.crop_widget and self.exam_image:
            base_name = os.path.splitext(self.exam_image)[0]
            config_file = base_name + "_config.json"

            self.crop_widget.save_regions_config(config_file)
            QMessageBox.information(self, "成功", f"配置已保存到：\n{config_file}")

    def load_config(self):
        """加载切割配置"""
        if self.crop_widget:
            filename, _ = QFileDialog.getOpenFileName(
                self, "选择配置文件", "", "JSON Files (*.json)"
            )

            if filename:
                if self.crop_widget.load_regions_config(filename):
                    QMessageBox.information(self, "成功", "配置加载成功")
                else:
                    QMessageBox.critical(self, "错误", "配置加载失败")

    def validate_data(self):
        """校验Excel数据"""
        result = ""

        # 检查是否有数据
        if self.df is None or len(self.df) == 0:
            return "❌ 数据为空"

        # 查找姓名列
        name_column = None
        for col in self.df.columns:
            if '姓名' in str(col):
                name_column = col
                break

        if name_column:
            result += f"✓ 姓名列: {name_column}\n"
        else:
            result += f"⚠ 未找到姓名列\n"

        # 查找题目列
        score_columns = []
        objective_cols = []  # 客观题
        subjective_cols = []  # 主观题

        for col in self.df.columns:
            col_str = str(col)
            if '满分' in col_str:
                score_columns.append(col)
                if '客-' in col_str:
                    objective_cols.append(col)
                elif '主-' in col_str:
                    subjective_cols.append(col)

        result += f"✓ 题目列总数: {len(score_columns)}\n"
        result += f"  - 客观题: {len(objective_cols)}\n"
        result += f"  - 主观题: {len(subjective_cols)}\n\n"

        # 检查数据类型
        if len(score_columns) > 0:
            # 抽样检查前3列
            sample_cols = score_columns[:min(3, len(score_columns))]

            for col in sample_cols:
                col_name = str(col)[:20]  # 截断长列名
                values = self.df[col].dropna()

                if len(values) == 0:
                    result += f"⚠ {col_name}: 全部为空\n"
                    continue

                # 检查数据类型
                is_numeric = True
                is_choice = True
                numeric_count = 0
                choice_count = 0
                invalid_count = 0

                for val in values[:10]:  # 检查前10个值
                    val_str = str(val).strip()

                    # 尝试转换为数字
                    try:
                        float(val)
                        numeric_count += 1
                    except:
                        is_numeric = False

                    # 检查是否为选项
                    if val_str in ['A', 'B', 'C', 'D', 'E', 'F']:
                        choice_count += 1
                    elif val_str not in ['-', '', 'nan']:
                        is_choice = False
                        if not val_str.replace('.', '').replace('-', '').isdigit():
                            invalid_count += 1

                if numeric_count > 0:
                    result += f"✓ {col_name}: 数字类型 ({numeric_count}/10)\n"
                elif choice_count > 0:
                    result += f"⚠ {col_name}: 选项类型 (需标准答案)\n"
                elif invalid_count > 0:
                    result += f"❌ {col_name}: 数据异常\n"

        # 检查是否所有主观题都是数字
        if len(subjective_cols) > 0:
            result += f"\n主观题检查:\n"
            numeric_subjective = 0
            for col in subjective_cols[:5]:  # 检查前5个主观题
                sample = self.df[col].dropna().iloc[0] if len(self.df[col].dropna()) > 0 else None
                if sample is not None:
                    try:
                        float(sample)
                        numeric_subjective += 1
                    except:
                        pass

            if numeric_subjective > 0:
                result += f"✓ 主观题包含数字分数 ({numeric_subjective}/{min(5, len(subjective_cols))})\n"
            else:
                result += f"❌ 主观题不是数字格式\n"

        # 检查客观题
        if len(objective_cols) > 0:
            result += f"\n客观题提示:\n"
            result += f"⚠ 检测到{len(objective_cols)}个客观题列\n"
            result += f"⚠ 客观题为选项格式(A/B/C/D)\n"
            result += f"⚠ 需要标准答案才能统计\n"
            result += f"💡 建议：只分析主观题列\n"

        return result

    def check_ready(self):
        """检查是否可以生成PPT"""
        print(f"检查是否可以生成PPT...")
        print(f"  Excel文件: {self.excel_file}")
        print(f"  试卷文件: {self.exam_image}")
        print(f"  crop_widget: {self.crop_widget is not None}")
        if self.crop_widget:
            print(f"  框选数量: {len(self.crop_widget.rectangles)}")

        # 不再禁用按钮，让用户始终可以点击并获得反馈
        if self.excel_file and self.exam_image and self.crop_widget:
            if len(self.crop_widget.rectangles) > 0:
                print("  >>> 所有条件满足，可以生成PPT")
                # 改变按钮样式表示就绪
                self.generate_btn.setStyleSheet(
                    "background-color: #4CAF50; color: white; font-size: 14px; padding: 10px; font-weight: bold;")
                self.generate_btn.setText("✓ 生成PPT")
            else:
                print("  >>> 框选数量为0")
                self.generate_btn.setStyleSheet(
                    "background-color: #9E9E9E; color: white; font-size: 14px; padding: 10px;")
                self.generate_btn.setText("生成PPT (请先框选题目)")
        else:
            print("  >>> 条件不满足")
            self.generate_btn.setStyleSheet("background-color: #9E9E9E; color: white; font-size: 14px; padding: 10px;")
            if not self.excel_file:
                self.generate_btn.setText("生成PPT (请先选择成绩单)")
            elif not self.exam_image:
                self.generate_btn.setText("生成PPT (请先选择试卷)")
            else:
                self.generate_btn.setText("生成PPT (请先框选题目)")

    def debug_generate(self):
        """调试用：显示当前状态"""
        print("\n" + "=" * 50)
        print("【调试信息】")
        print(f"Excel文件: {self.excel_file}")
        print(f"试卷文件: {self.exam_image}")
        print(f"DataFrame: {self.df is not None}")
        if self.df is not None:
            print(f"  DataFrame行数: {len(self.df)}")
            print(f"  DataFrame列数: {len(self.df.columns)}")
        print(f"crop_widget: {self.crop_widget is not None}")
        if self.crop_widget:
            print(f"  框选区域数量: {len(self.crop_widget.rectangles)}")
        print(f"生成按钮状态: {'启用' if self.generate_btn.isEnabled() else '禁用'}")
        print("=" * 50 + "\n")

        msg = f"Excel: {'有' if self.excel_file else '无'}\n"
        msg += f"试卷: {'有' if self.exam_image else '无'}\n"
        msg += f"DataFrame: {'有' if self.df is not None else '无'}\n"
        msg += f"框选区域: {len(self.crop_widget.rectangles) if self.crop_widget else 0}\n"
        msg += f"生成按钮: {'启用' if self.generate_btn.isEnabled() else '禁用'}"

        QMessageBox.information(self, "调试信息", msg)

        # 如果一切就绪，尝试生成
        if self.df is not None and self.crop_widget and len(self.crop_widget.rectangles) > 0:
            reply = QMessageBox.question(self, "确认", "条件满足，是否立即生成PPT？",
                                         QMessageBox.Yes | QMessageBox.No)
            if reply == QMessageBox.Yes:
                self.generate_ppt()

    def generate_ppt(self):
        """生成PPT"""
        print("=" * 50)
        print("点击了生成PPT按钮！")
        print(f"Excel文件: {self.excel_file}")
        print(f"试卷文件: {self.exam_image}")
        print(f"DataFrame是否为空: {self.df is None}")
        print(f"crop_widget是否存在: {self.crop_widget is not None}")
        if self.crop_widget:
            print(f"框选区域数量: {len(self.crop_widget.rectangles)}")

        # 检查所有必要条件
        missing = []
        if self.df is None:
            missing.append("• 请先选择并加载成绩单Excel文件")
        if not self.exam_image:
            missing.append("• 请先选择试卷图片或PDF文件")
        if not self.crop_widget or len(self.crop_widget.rectangles) == 0:
            missing.append("• 请在试卷图片上框选题目区域\n  （用鼠标拖动框选每道题）")

        if missing:
            msg = "缺少必要条件：\n\n" + "\n".join(missing)
            QMessageBox.warning(self, "无法生成PPT", msg)
            return

        try:
            print("开始处理...")
            # 获取切割的题目图片
            print("正在切割题目图片...")
            regions = self.crop_widget.get_cropped_regions()
            print(f"已切割 {len(regions)} 个区域")

            # 分析每道题的答题情况
            print("正在分析答题情况...")
            question_stats = self.analyze_questions()
            print(f"已分析 {len(question_stats)} 道题目")

            # 创建PPT（16:9宽屏）
            print("正在创建PPT（16:9宽屏）...")
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9宽屏
            prs.slide_height = Inches(7.5)

            # 添加标题页（宽屏版）
            print("创建标题页...")
            title_slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 标题页背景
            background = title_slide.shapes.add_shape(
                1, 0, 0, prs.slide_width, prs.slide_height
            )
            background.fill.solid()
            background.fill.fore_color.rgb = RGBColor(41, 128, 185)
            background.line.color.rgb = RGBColor(41, 128, 185)

            # 主标题
            title_box = title_slide.shapes.add_textbox(
                Inches(2), Inches(2.5), Inches(9.333), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_p = title_frame.paragraphs[0]
            title_p.text = "试卷分析报告"
            title_p.font.size = Pt(54)
            title_p.font.bold = True
            title_p.font.color.rgb = RGBColor(255, 255, 255)
            title_p.alignment = 1

            # 副标题
            subtitle_box = title_slide.shapes.add_textbox(
                Inches(2), Inches(4.2), Inches(9.333), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = f"共 {len(regions)} 道题目 | 学生人数: {len(self.df)}"
            subtitle_p.font.size = Pt(20)
            subtitle_p.font.color.rgb = RGBColor(236, 240, 241)
            subtitle_p.alignment = 1

            # 为每道题创建页面（宽屏版）
            for qnum, img in regions:
                print(f"正在创建第{qnum}题...")
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # 浅色背景
                bg_shape = slide.shapes.add_shape(
                    1, 0, 0, prs.slide_width, prs.slide_height
                )
                bg_shape.fill.solid()
                bg_shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
                bg_shape.line.fill.background()

                # 顶部蓝色条
                header_bar = slide.shapes.add_shape(
                    1, 0, 0, prs.slide_width, Inches(0.6)
                )
                header_bar.fill.solid()
                header_bar.fill.fore_color.rgb = RGBColor(52, 152, 219)
                header_bar.line.fill.background()

                # 题号
                title_box = slide.shapes.add_textbox(
                    Inches(0.3), Inches(0.05), Inches(2), Inches(0.5)
                )
                title_frame = title_box.text_frame
                title_p = title_frame.paragraphs[0]
                title_p.text = f"第 {qnum} 题"
                title_p.font.size = Pt(28)
                title_p.font.bold = True
                title_p.font.color.rgb = RGBColor(255, 255, 255)

                # 图片背景框（上半部分）
                img_bg = slide.shapes.add_shape(
                    1, Inches(0.4), Inches(0.8), Inches(12.5), Inches(4.0)
                )
                img_bg.fill.solid()
                img_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
                img_bg.line.color.rgb = RGBColor(189, 195, 199)
                img_bg.line.width = Pt(1)

                # 题目图片（自动适应框内，不超出）
                img_stream = io.BytesIO()
                img.save(img_stream, format='PNG')
                img_stream.seek(0)

                # 框的可用空间（留边距）
                frame_width = Inches(12.3)  # 12.5 - 0.2边距
                frame_height = Inches(3.8)  # 4.0 - 0.2边距
                frame_left = Inches(0.5)
                frame_top = Inches(0.9)

                # 获取图片原始尺寸
                img_width, img_height = img.size
                img_aspect = img_width / img_height
                frame_aspect = frame_width / frame_height

                # 计算缩放后的尺寸（保持比例，不超出框）
                if img_aspect > frame_aspect:
                    # 图片更宽，以宽度为准
                    pic_width = frame_width
                    pic_height = frame_width / img_aspect
                else:
                    # 图片更高，以高度为准
                    pic_height = frame_height
                    pic_width = frame_height * img_aspect

                # 居中显示
                pic_left = frame_left + (frame_width - pic_width) / 2
                pic_top = frame_top + (frame_height - pic_height) / 2

                pic = slide.shapes.add_picture(
                    img_stream, pic_left, pic_top,
                    width=pic_width, height=pic_height
                )
                print(f"  已添加题目图片 (尺寸: {pic_width / Inches(1):.1f}x{pic_height / Inches(1):.1f}英寸)")

                # 统计信息（下半部分）
                if qnum in question_stats:
                    stats = question_stats[qnum]
                    print(f"  第{qnum}题统计: 答对{stats['correct_count']}人, 答错{stats['wrong_count']}人")

                    # 信息背景框（下半部分）
                    info_bg = slide.shapes.add_shape(
                        1, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.2)
                    )
                    info_bg.fill.solid()
                    info_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    info_bg.line.color.rgb = RGBColor(189, 195, 199)
                    info_bg.line.width = Pt(1)

                    # 统计信息条（彩色）
                    correct_rate = stats['correct_rate']
                    if correct_rate >= 80:
                        stat_color = RGBColor(46, 204, 113)
                    elif correct_rate >= 60:
                        stat_color = RGBColor(241, 196, 15)
                    else:
                        stat_color = RGBColor(231, 76, 60)

                    stat_bar = slide.shapes.add_shape(
                        1, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.5)
                    )
                    stat_bar.fill.solid()
                    stat_bar.fill.fore_color.rgb = stat_color
                    stat_bar.line.fill.background()

                    # 统计数字（在彩色条上）
                    stat_text_box = slide.shapes.add_textbox(
                        Inches(0.6), Inches(5.05), Inches(12.1), Inches(0.4)
                    )
                    stat_tf = stat_text_box.text_frame
                    stat_p = stat_tf.paragraphs[0]
                    stat_p.text = f"正确率: {correct_rate:.1f}%  |  答对: {stats['correct_count']}人  |  答错: {stats['wrong_count']}人"
                    stat_p.font.size = Pt(18)
                    stat_p.font.bold = True
                    stat_p.font.color.rgb = RGBColor(255, 255, 255)

                    # 答错学生标题
                    title_box = slide.shapes.add_textbox(
                        Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.4)
                    )
                    title_tf = title_box.text_frame
                    title_p = title_tf.paragraphs[0]
                    title_p.text = "答错学生名单："
                    title_p.font.size = Pt(18)
                    title_p.font.bold = True
                    title_p.font.color.rgb = RGBColor(52, 73, 94)

                    # 学生列表
                    wrong_students = stats['wrong_students']
                    student_box = slide.shapes.add_textbox(
                        Inches(0.6), Inches(6.1), Inches(12.1), Inches(1.0)
                    )
                    student_tf = student_box.text_frame
                    student_tf.word_wrap = True
                    student_p = student_tf.paragraphs[0]

                    if wrong_students:
                        student_p.text = "、".join(wrong_students)
                        student_p.font.size = Pt(16)
                        student_p.font.color.rgb = RGBColor(231, 76, 60)
                        student_p.line_spacing = 1.2
                        print(f"  答错学生: {len(wrong_students)}人")
                    else:
                        student_p.text = "🎉 全部答对！"
                        student_p.font.size = Pt(20)
                        student_p.font.bold = True
                        student_p.font.color.rgb = RGBColor(46, 204, 113)
                        student_p.alignment = 1
                        print(f"  全部答对！")
                else:
                    print(f"  警告: 第{qnum}题没有找到统计数据！")
                    # 下方信息框
                    info_bg = slide.shapes.add_shape(
                        1, Inches(0.4), Inches(5.0), Inches(12.5), Inches(1.2)
                    )
                    info_bg.fill.solid()
                    info_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    info_bg.line.color.rgb = RGBColor(189, 195, 199)
                    info_bg.line.width = Pt(1)

                    textbox = slide.shapes.add_textbox(
                        Inches(0.6), Inches(5.3), Inches(12.1), Inches(0.8)
                    )
                    text_frame = textbox.text_frame
                    p = text_frame.paragraphs[0]
                    p.text = "（暂无统计数据）"
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(149, 165, 166)
                    p.alignment = 1

                # 页脚（页码）
                footer_box = slide.shapes.add_textbox(
                    Inches(12.5), Inches(7.2), Inches(0.7), Inches(0.25)
                )
                footer_tf = footer_box.text_frame
                footer_p = footer_tf.paragraphs[0]
                footer_p.text = f"{qnum}/{len(regions)}"
                footer_p.font.size = Pt(11)
                footer_p.font.color.rgb = RGBColor(127, 140, 141)
                footer_p.alignment = 2

            # 保存PPT
            output_file = "试卷分析报告.pptx"
            print(f"正在保存PPT到: {output_file}")
            prs.save(output_file)
            print("PPT保存成功！")
            print("=" * 50)

            QMessageBox.information(self, "成功",
                                    f"PPT已生成！\n保存为：{output_file}\n"
                                    f"共 {len(regions)} 页")

        except Exception as e:
            print(f"错误: {e}")
            QMessageBox.critical(self, "错误", f"生成PPT失败：\n{str(e)}")
            import traceback
            traceback.print_exc()

    def analyze_questions(self):
        """分析每道题的答题情况"""
        stats = {}

        print("\n" + "=" * 80)
        print("【开始分析题目数据】")
        print("=" * 80)

        # 打印DataFrame基本信息
        print(f"\nDataFrame形状: {self.df.shape}")
        print(f"DataFrame列名（前10个）: {self.df.columns.tolist()[:10]}")

        # 查找成绩列和姓名列
        score_columns = []
        name_column = None

        # 寻找姓名列
        for col in self.df.columns:
            col_str = str(col)
            if '姓名' in col_str or 'name' in col_str.lower() or '学生' in col_str:
                name_column = col
                break

        if name_column is None:
            # 尝试查找第二列或第三列作为姓名列
            if len(self.df.columns) > 2:
                name_column = self.df.columns[2]
            elif len(self.df.columns) > 1:
                name_column = self.df.columns[1]
            else:
                name_column = self.df.columns[0]

        print(f"\n姓名列: {name_column}")
        print(f"姓名列的前5个值: {self.df[name_column].head().tolist()}")

        # 寻找题目分数列
        # 格式可能是："主-1 (满分: 2)"、"第1题"、"1题"、"客-1 (满分: 1)"等
        all_question_columns = []
        subjective_columns = []  # 主观题列（主-X）
        objective_columns = []  # 客观题列（客-X）

        for col in self.df.columns:
            col_str = str(col)
            # 匹配包含"满分"、"题"的列，或纯数字列
            if ('满分' in col_str or '题' in col_str or
                    (col_str.replace('.', '').replace('-', '').isdigit() and col_str != name_column)):
                all_question_columns.append(col)

                # 区分主观题和客观题
                if '主-' in col_str or '主' in col_str[:2]:
                    subjective_columns.append(col)
                elif '客-' in col_str or '客' in col_str[:2]:
                    objective_columns.append(col)
                else:
                    # 默认按主观题处理
                    subjective_columns.append(col)

        print(f"\n发现题目列:")
        print(f"  总计: {len(all_question_columns)} 个")
        print(f"  主观题(主-X): {len(subjective_columns)} 个")
        print(f"  客观题(客-X): {len(objective_columns)} 个")

        # 优先使用主观题列（有实际分数）
        if len(subjective_columns) > 0:
            score_columns = subjective_columns
            print(f"\n✓ 使用主观题列进行分析（共 {len(score_columns)} 题）")
            print(f"主观题列（前5个）: {score_columns[:5]}")

            # 验证主观题是否真的是数字
            sample_col = score_columns[0]
            sample_values = self.df[sample_col].dropna()
            if len(sample_values) > 0:
                sample_val = sample_values.iloc[0]
                print(f"\n数据类型验证:")
                print(f"  示例列: {sample_col}")
                print(f"  示例值: {sample_val} (类型: {type(sample_val).__name__})")
                try:
                    float(sample_val)
                    print(f"  ✓ 可转换为数字")
                except:
                    print(f"  ❌ 无法转换为数字！")
                    print(f"  ⚠ 警告：主观题列包含非数字数据")
        else:
            # 如果没有主观题，尝试所有题目列
            score_columns = all_question_columns
            print(f"\n⚠ 未找到主观题列，尝试使用所有题目列（共 {len(score_columns)} 题）")
            print(f"题目列（前5个）: {score_columns[:5]}")
            print(f"⚠ 注意：可能包含选项类型数据(A/B/C/D)，无法统计对错")

        # 分析每道题
        print("\n开始逐题分析...")
        for idx, col in enumerate(score_columns, 1):
            if idx <= 3:  # 只详细打印前3题
                print(f"\n--- 分析第{idx}题 (列名: {col}) ---")

            wrong_students = []
            correct_count = 0
            wrong_count = 0
            total_students = 0

            # 打印这一列的前5个值
            if idx <= 3:
                print(f"该列前5个值: {self.df[col].head().tolist()}")

            row_num = 0
            for _, row in self.df.iterrows():
                row_num += 1

                # 跳过空行或标题行
                if pd.isna(row[name_column]) or str(row[name_column]).strip() == '':
                    if idx <= 3 and row_num <= 3:
                        print(f"  行{row_num}: 跳过空姓名")
                    continue

                student_name = str(row[name_column]).strip()

                # 跳过"姓名"、"学生"等标题行
                if student_name in ['姓名', '学生', 'name', 'Name']:
                    if idx <= 3 and row_num <= 3:
                        print(f"  行{row_num}: 跳过标题行 '{student_name}'")
                    continue

                score = row[col]
                score_str = str(score).strip()

                # 跳过未作答的(-、空值等)
                if score_str in ['-', '', 'nan', 'NaN']:
                    if idx <= 3 and total_students <= 5:
                        print(f"  学生 '{student_name}': 未作答")
                    continue

                # 尝试判断是选项还是分数
                try:
                    # 先尝试转换为数字（分数类型）
                    score_val = float(score)
                    total_students += 1

                    if idx <= 3 and total_students <= 5:
                        print(f"  学生 '{student_name}': 得分={score_val}")

                    # 如果得分为0，认为答错
                    if score_val == 0 or score_val == 0.0:
                        wrong_students.append(student_name)
                        wrong_count += 1
                    else:
                        correct_count += 1

                except (ValueError, TypeError):
                    # 如果无法转换为数字，可能是选项类型（A/B/C/D）
                    # 对于选项类型，我们暂时无法判断对错（需要标准答案）
                    # 这里先统计答案分布，将少数答案视为错误
                    if idx <= 3 and total_students <= 5:
                        print(f"  学生 '{student_name}': 选择={score_str} (选项类型，暂无法判断对错)")
                    # 跳过选项类型的处理
                    continue

            total = correct_count + wrong_count
            correct_rate = (correct_count / total * 100) if total > 0 else 0

            if idx <= 3:
                print(
                    f"第{idx}题统计: 总人数={total}, 答对={correct_count}, 答错={wrong_count}, 正确率={correct_rate:.1f}%")
                print(f"答错学生: {wrong_students[:5]}")

            stats[idx] = {
                'wrong_students': wrong_students,
                'correct_count': correct_count,
                'wrong_count': wrong_count,
                'correct_rate': correct_rate,
                'column_name': col
            }

        print("\n" + "=" * 80)
        print(f"【分析完成】共分析 {len(stats)} 道题")
        print("=" * 80 + "\n")

        return stats


def main():
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())


if __name__ == '__main__':
    main()
